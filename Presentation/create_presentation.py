#!/usr/bin/env python3
"""
Maykana ERP - PowerPoint Presentation Generator
Creates a professional bilingual presentation about the ERP system
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Brand Colors from the project
PRIMARY_COLOR = RGBColor(9, 55, 56)  # #093738
PRIMARY_HOVER = RGBColor(10, 72, 73)  # #0a4849
SUCCESS_COLOR = RGBColor(44, 194, 141)  # #2cc28d
BG_LIGHT = RGBColor(241, 245, 249)  # #F1F5F9
BORDER_COLOR = RGBColor(226, 226, 226)  # #e2e2e2
ACCENT_COLOR = RGBColor(100, 200, 150)  # Custom accent

def create_presentation():
    """Create the main presentation with all slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs)
    
    # Slide 2: Overview
    add_overview_slide(prs)
    
    # Slide 3: Goals & Benefits
    add_goals_benefits_slide(prs)
    
    # Slide 4: Technical Stack
    add_technical_stack_slide(prs)
    
    # Slide 5-14: Modules
    add_accounting_module_slide(prs)
    add_purchases_module_slide(prs)
    add_sales_module_slide(prs)
    add_warehouses_module_slide(prs)
    add_hr_module_slide(prs)
    add_assets_module_slide(prs)
    add_competitions_module_slide(prs)
    add_strategy_module_slide(prs)
    add_workflow_slide(prs)
    add_reports_slide(prs)
    
    # Slide 15: UI/UX Design
    add_ui_design_slide(prs)
    
    # Slide 16: Workflow Engine
    add_workflow_engine_detail_slide(prs)
    
    # Slide 17: Integration
    add_integration_slide(prs)
    
    # Slide 18: Security
    add_security_slide(prs)
    
    # Slide 19: Reports & Analytics
    add_analytics_slide(prs)
    
    # Slide 20: Roadmap & Next Steps
    add_roadmap_slide(prs)
    
    # Save presentation
    output_path = 'Maykana_ERP_Presentation.pptx'
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    return output_path

def add_blank_slide(prs):
    """Add a blank slide layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide

def add_title_shape(slide, title_ar, title_en, top=0.5):
    """Add bilingual title to slide"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9), Inches(0.8)
    )
    text_frame = title_box.text_frame
    
    # Arabic title
    p = text_frame.paragraphs[0]
    p.text = title_ar
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.font.name = 'IBM Plex Sans Arabic'
    
    # English title
    p = text_frame.add_paragraph()
    p.text = title_en
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY_HOVER
    p.font.name = 'Calibri'
    
    return title_box

def add_content_box(slide, content, left, top, width, height, is_rtl=True):
    """Add content box with proper formatting"""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = box.text_frame
    text_frame.word_wrap = True
    
    paragraphs = content if isinstance(content, list) else [content]
    
    for i, line in enumerate(paragraphs):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
        
    return box

def add_bullet_list(slide, items, left, top, width, height, is_rtl=True):
    """Add bullet list with bilingual support"""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = box.text_frame
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)
    
    return box

def add_background_shape(slide, color=BG_LIGHT):
    """Add background rectangle"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

# ==================== SLIDE FUNCTIONS ====================

def add_title_slide(prs):
    """Slide 1: Title/Cover"""
    slide = add_blank_slide(prs)
    
    # Background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.color.rgb = PRIMARY_COLOR
    
    # Main Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "نظام ميكنة لتخطيط موارد المؤسسات"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # English subtitle
    p = tf.add_paragraph()
    p.text = "Maykana ERP System"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.color.rgb = SUCCESS_COLOR
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.5), Inches(6), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "نظام متكامل لإدارة كافة عمليات المؤسسة"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "Complete Enterprise Resource Planning Solution"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = BG_LIGHT

def add_overview_slide(prs):
    """Slide 2: System Overview"""
    slide = add_blank_slide(prs)
    
    # White background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    # Title
    add_title_shape(slide, "نظرة عامة", "System Overview")
    
    # Content
    content = [
        "• نظام ERP متكامل مصمم للشركات السعودية والعربية",
        "  Complete ERP system designed for Saudi and Arabic companies",
        "",
        "• واجهة عربية أصلية مع دعم كامل للغة العربية (RTL)",
        "  Native Arabic interface with full RTL support",
        "",
        "• معمارية Monorepo حديثة باستخدام أحدث التقنيات",
        "  Modern Monorepo architecture using latest technologies",
        "",
        "• نظام شامل يغطي 10+ وحدات رئيسية",
        "  Comprehensive system covering 10+ main modules",
        "",
        "• تكامل كامل بين جميع الوحدات",
        "  Full integration between all modules"
    ]
    
    add_content_box(slide, content, 1, 2, 8, 4.5)

def add_goals_benefits_slide(prs):
    """Slide 3: Goals & Benefits"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "الأهداف والفوائد", "Goals & Benefits")
    
    # Right column - Goals
    goals = [
        "الأهداف | Goals:",
        "• أتمتة جميع العمليات الإدارية",
        "• تحسين الكفاءة التشغيلية",
        "• توحيد البيانات في نظام واحد",
        "• تقليل الأخطاء البشرية",
        "• تسريع عملية اتخاذ القرار"
    ]
    add_bullet_list(slide, goals, 5.5, 2, 4, 4.5)
    
    # Left column - Benefits
    benefits = [
        "الفوائد | Benefits:",
        "• توفير الوقت والجهد",
        "• تقليل التكاليف التشغيلية",
        "• تحسين دقة البيانات",
        "• تقارير فورية شاملة",
        "• سهولة التدقيق والمراجعة",
        "• الامتثال للمعايير المحاسبية"
    ]
    add_bullet_list(slide, benefits, 0.5, 2, 4.5, 4.5)

def add_technical_stack_slide(prs):
    """Slide 4: Technical Stack"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "التقنيات المستخدمة", "Technology Stack")
    
    # Frontend
    frontend = [
        "واجهة المستخدم | Frontend:",
        "• React 18 - مكتبة بناء الواجهات",
        "• TypeScript 5.8 - لغة البرمجة",
        "• Vite - أداة التطوير السريعة",
        "• Tailwind CSS - إطار التصميم",
        "• Shadcn/ui - مكتبة المكونات"
    ]
    add_bullet_list(slide, frontend, 0.5, 2, 4.5, 2.5)
    
    # Backend Architecture
    backend = [
        "البنية التحتية | Architecture:",
        "• Turborepo - نظام Monorepo",
        "• pnpm - مدير الحزم",
        "• Redux Toolkit - إدارة الحالة",
        "• React Router v6 - التنقل",
        "• IBM Plex Sans Arabic - الخط العربي"
    ]
    add_bullet_list(slide, backend, 5.5, 2, 4, 2.5)
    
    # Features
    features = [
        "المميزات التقنية | Technical Features:",
        "• معمارية Monorepo للمشاركة الفعالة",
        "• أداء عالي مع Lazy Loading",
        "• تصميم متجاوب لجميع الأجهزة",
        "• دعم كامل للغة العربية (RTL)",
        "• نظام مكونات قابل لإعادة الاستخدام"
    ]
    add_bullet_list(slide, features, 0.5, 4.8, 9, 2.5)

def add_accounting_module_slide(prs):
    """Slide 5: Accounting Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "🧮 وحدة إدارة الحسابات", "Accounting Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• القيود المحاسبية - تسجيل وإدارة القيود",
        "• سندات القبض والصرف - إدارة المستندات المالية",
        "• عهد نقدية - تتبع العهد المالية",
        "• شجرة الحسابات - دليل حسابات شامل",
        "• التقارير المالية - ميزان المراجعة، قائمة الدخل، المركز المالي",
        "• مراكز التكلفة - تخصيص النفقات",
        "• العملات المتعددة - دعم التعاملات الدولية",
        "",
        "Main Functions:",
        "Accounting Entries • Receipt & Payment Vouchers • Cash Custody",
        "Chart of Accounts • Financial Reports • Cost Centers • Multi-Currency"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_purchases_module_slide(prs):
    """Slide 6: Purchases Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "🛒 وحدة إدارة المشتريات", "Purchases Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• إدارة الموردين - قاعدة بيانات شاملة للموردين",
        "• طلبات الشراء - إنشاء ومتابعة طلبات المواد",
        "• عروض الأسعار - طلب ومقارنة عروض الموردين",
        "• أوامر الشراء - إصدار وتتبع أوامر الشراء",
        "• استلام المواد - تسجيل استلام الطلبات",
        "• فواتير المشتريات - معالجة فواتير الموردين",
        "• التكامل مع المخازن والحسابات",
        "",
        "Main Functions:",
        "Supplier Management • Purchase Requests • Price Quotes",
        "Purchase Orders • Material Receipt • Purchase Invoices"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_sales_module_slide(prs):
    """Slide 7: Sales Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "💼 وحدة إدارة المبيعات", "Sales Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• إدارة العملاء - قاعدة بيانات العملاء الشاملة",
        "• عروض الأسعار - إنشاء عروض احترافية للعملاء",
        "• أوامر البيع - إدارة طلبات العملاء",
        "• فواتير المبيعات - إصدار ومتابعة الفواتير",
        "• إذن تسليم - تسجيل عمليات التسليم",
        "• قوائم الأسعار - إدارة الأسعار المختلفة",
        "• مندوبي المبيعات - متابعة الأداء والعمولات",
        "",
        "Main Functions:",
        "Customer Management • Quotations • Sales Orders • Invoices",
        "Delivery Notes • Price Lists • Sales Representatives"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_warehouses_module_slide(prs):
    """Slide 8: Warehouses Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "📦 وحدة إدارة المستودعات", "Warehouses Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• مواد المخزون - تسجيل وإدارة جميع المواد",
        "• حركات المخزون - تتبع جميع حركات الصرف والإضافة",
        "• جرد المخزون - إجراء الجرد الدوري والمفاجئ",
        "• طلبات الانفاق - إدارة طلبات صرف المواد",
        "• نقل المواد - نقل بين المستودعات",
        "• الأرصدة الافتتاحية - تسجيل الأرصدة الأولية",
        "• التقييم - طرق تقييم مختلفة (FIFO, LIFO, Average)",
        "",
        "Main Functions:",
        "Inventory Materials • Stock Movements • Inventory Count",
        "Material Requests • Transfers • Opening Balances • Valuation"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_hr_module_slide(prs):
    """Slide 9: HR Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "👥 وحدة الموارد البشرية", "Human Resources Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• ملفات الموظفين - إدارة بيانات الموظفين الكاملة",
        "• الحضور والانصراف - تتبع أوقات العمل",
        "• الإجازات - إدارة طلبات واستحقاقات الإجازات",
        "• الرواتب - حساب وصرف الرواتب الشهرية",
        "• التقييم والتطوير - تقييم الأداء والتدريب",
        "• التوظيف - إدارة عملية التوظيف من البداية للنهاية",
        "• العمل عن بعد - إدارة سياسات وطلبات العمل عن بعد",
        "",
        "Main Functions:",
        "Employee Records • Attendance • Leave Management • Payroll",
        "Performance • Recruitment • Remote Work"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_assets_module_slide(prs):
    """Slide 10: Assets Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "🏢 وحدة إدارة الأصول", "Assets Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• سجل الأصول - قاعدة بيانات شاملة لجميع الأصول",
        "• حركات الأصول - نقل وتحويل الأصول",
        "• الاستهلاك - حساب الاستهلاك تلقائياً بطرق مختلفة",
        "• الصيانة - جدولة ومتابعة أعمال الصيانة",
        "• التقييم - إعادة تقييم قيمة الأصول",
        "• البيع والاستبعاد - إدارة عمليات التخلص من الأصول",
        "• التقارير - تقارير شاملة عن حالة الأصول",
        "",
        "Main Functions:",
        "Asset Register • Movements • Depreciation • Maintenance",
        "Revaluation • Disposal • Comprehensive Reports"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_competitions_module_slide(prs):
    """Slide 11: Competitions Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "🏆 وحدة إدارة المنافسات", "Competitions Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• تأهيل الموردين - تسجيل وتقييم الموردين المؤهلين",
        "• تشكيل اللجان - إنشاء لجان التقييم",
        "• معايير التقييم - تحديد معايير تقييم العروض",
        "• إطلاق المنافسة - نشر والإعلان عن المنافسات",
        "• استقبال وفتح العروض - إدارة عملية استلام العروض",
        "• التقييم والترسية - تقييم العروض واختيار الفائز",
        "• العقود والاتفاقيات - إدارة التعاقدات",
        "• الضمانات البنكية - متابعة الضمانات",
        "",
        "Main Functions:",
        "Vendor Qualification • Committee Formation • Evaluation Criteria",
        "Competition Launch • Offers Management • Award • Contracts"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_strategy_module_slide(prs):
    """Slide 12: Strategy Module"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "🎯 وحدة الإستراتيجية", "Strategy Management Module")
    
    features = [
        "الوظائف الرئيسية:",
        "• الخطط الاستراتيجية - وضع ومتابعة الخطط",
        "• المشاريع - إدارة المشاريع الاستراتيجية",
        "• المهام - تعيين ومتابعة المهام",
        "• الاجتماعات - جدولة اجتماعات الإدارة",
        "• الوثائق - إدارة المستندات الاستراتيجية",
        "• تتبع الأداء - متابعة تحقيق الأهداف",
        "• المؤشرات (KPIs) - قياس الأداء المؤسسي",
        "",
        "Main Functions:",
        "Strategic Plans • Projects • Tasks • Meetings",
        "Documents • Performance Tracking • KPIs"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_workflow_slide(prs):
    """Slide 13: Workflow Engine"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "⚙️ محرك سير العمل", "Workflow Engine")
    
    features = [
        "الوظائف الرئيسية:",
        "• مسارات العمل - تصميم مسارات الموافقات المخصصة",
        "• قوائم التحقق - إنشاء قوالب التحقق والمراجعة",
        "• الأدوار والصلاحيات - تحديد المسؤوليات",
        "• الإشعارات التلقائية - تنبيه المسؤولين",
        "• تتبع المهام - متابعة حالة الطلبات",
        "• التكامل الشامل - ربط مع جميع الوحدات",
        "• السجلات - حفظ سجل كامل للإجراءات",
        "",
        "Main Functions:",
        "Custom Workflows • Verification Templates • Roles & Permissions",
        "Auto Notifications • Task Tracking • Full Integration • Audit Trail"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_reports_slide(prs):
    """Slide 14: Reports"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "📊 نظام التقارير", "Reports System")
    
    features = [
        "التقارير المتوفرة:",
        "• تقارير المحاسبة - ميزان المراجعة، قائمة الدخل، المركز المالي",
        "• تقارير المشتريات - تحليل المشتريات، مقارنة الموردين",
        "• تقارير المبيعات - تحليل المبيعات، أداء المندوبين",
        "• تقارير المخازن - حالة المخزون، حركة المخزون",
        "• تقارير الموارد البشرية - الحضور، الرواتب، الأداء",
        "• تقارير الأصول - حالة الأصول، الاستهلاك",
        "• تقارير مخصصة - تصميم تقارير حسب الحاجة",
        "",
        "Available Reports:",
        "Accounting • Purchases • Sales • Inventory • HR • Assets • Custom Reports"
    ]
    add_bullet_list(slide, features, 1, 2, 8, 5)

def add_ui_design_slide(prs):
    """Slide 15: UI/UX Design"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "🎨 التصميم وتجربة المستخدم", "UI/UX Design")
    
    design_features = [
        "مميزات التصميم:",
        "• واجهة عربية أصلية 100%",
        "  100% Native Arabic Interface",
        "",
        "• دعم كامل للاتجاه من اليمين لليسار (RTL)",
        "  Full Right-to-Left (RTL) Support",
        "",
        "• تصميم متجاوب لجميع الأجهزة",
        "  Responsive Design for All Devices",
        "",
        "• ألوان احترافية (#093738 كلون رئيسي)",
        "  Professional Color Scheme",
        "",
        "• خط IBM Plex Sans Arabic للنصوص العربية",
        "  IBM Plex Arabic Font",
        "",
        "• واجهة بديهية سهلة الاستخدام",
        "  Intuitive User-Friendly Interface"
    ]
    add_content_box(slide, design_features, 1, 2, 8, 4.5)

def add_workflow_engine_detail_slide(prs):
    """Slide 16: Workflow Engine Details"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "⚙️ محرك سير العمل - التفاصيل", "Workflow Engine - Details")
    
    features = [
        "كيف يعمل محرك سير العمل:",
        "",
        "1️⃣ تصميم المسار - إنشاء مسار موافقات مخصص",
        "   Design workflow with custom approval paths",
        "",
        "2️⃣ تحديد المسؤولين - ربط كل خطوة بمسؤول محدد",
        "   Assign responsible users for each step",
        "",
        "3️⃣ الإشعارات التلقائية - إرسال تنبيهات للمسؤولين",
        "   Automatic notifications to stakeholders",
        "",
        "4️⃣ التتبع المباشر - متابعة حالة الطلب في كل مرحلة",
        "   Real-time request tracking",
        "",
        "5️⃣ السجل الكامل - حفظ جميع الإجراءات والتعليقات",
        "   Complete audit trail and comments"
    ]
    add_content_box(slide, features, 1, 2, 8, 5)

def add_integration_slide(prs):
    """Slide 17: System Integration"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "🔗 التكامل بين الأنظمة", "System Integration")
    
    integration = [
        "التكامل الشامل بين الوحدات:",
        "",
        "• المشتريات ← المخازن ← المحاسبة",
        "  Purchases → Warehouses → Accounting",
        "",
        "• المبيعات ← المخازن ← المحاسبة",
        "  Sales → Warehouses → Accounting",
        "",
        "• الموارد البشرية ← المحاسبة (الرواتب)",
        "  HR → Accounting (Payroll)",
        "",
        "• الأصول ← المحاسبة (الاستهلاك)",
        "  Assets → Accounting (Depreciation)",
        "",
        "• المنافسات ← المشتريات ← العقود",
        "  Competitions → Purchases → Contracts",
        "",
        "• جميع الوحدات ← محرك سير العمل",
        "  All Modules → Workflow Engine"
    ]
    add_content_box(slide, integration, 1, 2, 8, 5)

def add_security_slide(prs):
    """Slide 18: Security & Permissions"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT
    
    add_title_shape(slide, "🔐 الأمان والصلاحيات", "Security & Permissions")
    
    security = [
        "مميزات الأمان:",
        "• نظام صلاحيات متعدد المستويات",
        "  Multi-level permissions system",
        "",
        "• تسجيل دخول آمن مع المصادقة",
        "  Secure login with authentication",
        "",
        "• التحكم في الوصول حسب الدور",
        "  Role-based access control (RBAC)",
        "",
        "• سجل كامل للإجراءات (Audit Trail)",
        "  Complete audit trail",
        "",
        "• تشفير البيانات الحساسة",
        "  Data encryption for sensitive information",
        "",
        "• نسخ احتياطي تلقائي",
        "  Automatic backup system"
    ]
    add_content_box(slide, security, 1, 2, 8, 5)

def add_analytics_slide(prs):
    """Slide 19: Analytics & Insights"""
    slide = add_blank_slide(prs)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.color.rgb = RGBColor(255, 255, 255)
    
    add_title_shape(slide, "📈 التحليلات والرؤى", "Analytics & Insights")
    
    analytics = [
        "لوحات المعلومات والتحليلات:",
        "",
        "• لوحة معلومات تنفيذية شاملة",
        "  Comprehensive executive dashboard",
        "",
        "• مؤشرات الأداء الرئيسية (KPIs)",
        "  Key Performance Indicators",
        "",
        "• تحليل الاتجاهات والأنماط",
        "  Trend and pattern analysis",
        "",
        "• تقارير مصورة (Charts & Graphs)",
        "  Visual reports with charts",
        "",
        "• تصدير التقارير (PDF, Excel)",
        "  Export reports in multiple formats",
        "",
        "• تقارير مجدولة تلقائياً",
        "  Scheduled automated reports"
    ]
    add_content_box(slide, analytics, 1, 2, 8, 5)

def add_roadmap_slide(prs):
    """Slide 20: Roadmap & Next Steps"""
    slide = add_blank_slide(prs)
    
    # Gradient background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "خارطة الطريق والخطوات القادمة"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "Roadmap & Next Steps"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.color.rgb = SUCCESS_COLOR
    
    # Roadmap content
    roadmap_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(4)
    )
    tf = roadmap_box.text_frame
    
    phases = [
        "المرحلة الحالية:",
        "✅ البنية الأساسية والتصميم",
        "✅ وحدة الحسابات (جاري التطوير)",
        "",
        "المرحلة القادمة:",
        "🔄 إكمال الوحدات الأساسية (مشتريات، مبيعات، مخازن)",
        "🔄 تطوير محرك سير العمل",
        "🔄 نظام التقارير المتقدم",
        "",
        "المستقبل:",
        "🔜 تطبيق الجوال",
        "🔜 الذكاء الاصطناعي والتحليلات المتقدمة",
        "🔜 التكامل مع الأنظمة الخارجية"
    ]
    
    for i, line in enumerate(phases):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)
        if "المرحلة" in line or "Current" in line or "Next" in line or "Future" in line:
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = SUCCESS_COLOR

if __name__ == "__main__":
    create_presentation()
